import collections 
import collections.abc
from pptx import Presentation
import requests
import os
import time
import json
import sys

from ...tool import Tool

def build_tool(config) -> Tool:
    tool = Tool(
        "Slides Making",
        "This tool allows you to create ppt slides with text, paragraph, images, with good looking styles",
        name_for_model="Slides Making",
        description_for_model="This tool allows you to create ppt slides with text, paragraph, images, with good looking styles",
        logo_url="https://your-app-url.com/.well-known/logo.png",
        contact_email="bokesyo2000@gmail.com",
        legal_info_url="hello@legal.com"
    )

    CWD = os.getcwd() # path of current working directory
    LIB_DIR = os.path.dirname(__file__) # path of library
    TEMPLATE_DIR = os.path.join(LIB_DIR, "templates") # path of templates
    CACHE_DIR = os.path.join(CWD, "cache") # path of cache_dir
    IMAGE_BED_PATTERN = 'https://source.unsplash.com/featured/?{}' # url pattern for image bed

    if not os.path.exists(CACHE_DIR):
        os.makedirs(CACHE_DIR)
        # print(f"[system]: office.ppt: created cache directory: {CACHE_DIR}")

    # print(f"[system]: office.ppt_functional: TEMPLATE_DIR = {TEMPLATE_DIR}")
    # print(f"[system]: office.ppt_functional: CACHE_DIR = {CACHE_DIR}")

    ppt_file = None # a pointer to the powerpoint object

    def _return_timestamp():
        return str(time.time())

    def runtime_update_docstring(new_docstring: str) -> callable:
        """ This is a decorator that can help update the docstring at runtime
        """
        def decorator(func: callable) -> callable:
            func.__doc__ = new_docstring
            return func
        return decorator

    # Update the template list, then update the docstring of create_file
    ppt_template_names = []
    all_files = os.listdir(TEMPLATE_DIR)
    for file_name in all_files:
        if file_name.lower().endswith('.pptx'):
            ppt_template_names.append(file_name.split(".")[0])
    updated_docstring_create_file = f"""create_file(theme:str) -> None: Create a pptx file with specific theme, available thems: {' / '.join(ppt_template_names)}."""

    @tool.get("/create_file")
    @runtime_update_docstring(updated_docstring_create_file)
    def create_file(theme:str) -> str:
        """create_file(theme:str) -> None Create a pptx file with specific themes. Available themes: <update at runtime>
        """
        nonlocal ppt_file
        ppt_file = Presentation(os.path.join(TEMPLATE_DIR, f"{theme}.pptx"))
        return "created a ppt file."

    @tool.get("/get_image")
    def get_image(keywords:str) -> str:
        """get_image(keywords:str) -> str Get an image given comma separated keywords, return the image path.
        """
        picture_url = IMAGE_BED_PATTERN.format(keywords)
        response = requests.get(picture_url)
        img_local_path = os.path.join(CACHE_DIR, f"{_return_timestamp()}.jpg")
        with open(img_local_path, 'wb') as f:
            f.write(response.content)
        return img_local_path

    @tool.get("/add_first_page")
    def add_first_page(title:str, subtitle:str) -> str:
        """add_first_page(title:str, subtitle:str) -> None: Add the first page of ppt.
        """
        nonlocal ppt_file
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[0]) # layout for first page (title and subtitle only)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle
        
        return "added page"

    @tool.get("/add_text_page")
    def add_text_page(title:str, bullet_items:str) -> str:
        """add_text_page(title:str, bullet_items:str) -> None: Add text page (outline page is also applied).
        bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.
        """
        nonlocal ppt_file
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        tf = body_shape.text_frame

        bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1
        
        return "added page"
        
    @tool.get("/add_text_image_page")
    def add_text_image_page(title:str, bullet_items:str, image:str) -> str:
        """add_text_image_page(title:str, bullet_items:str, image:str) -> None: Add a text page with one image. (image should be a path)
        bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.
        """
        nonlocal ppt_file
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[3])
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.placeholders[1]
        
        tf = body_shape.text_frame
        bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1
        
        image_shape = slide.placeholders[2]
        slide.shapes.add_picture(image, image_shape.left, image_shape.top, image_shape.width, image_shape.height)
        
        return "added page"

    @tool.get("/submit_file")
    def submit_file() -> None:
        """submit_file() -> None: When all steps done, YOU MUST use submit_file() to submit your work.
        """
        nonlocal ppt_file
        file_path = os.path.join(CACHE_DIR, f"{_return_timestamp()}.pptx")
        ppt_file.save(file_path)
        # retreival_url = upload_file(file_path)

        return f"submitted. view ppt at {file_path}"

    return tool
